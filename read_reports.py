# -*- coding: utf-8 -*-
import openpyxl
from docx import Document
from pptx import Presentation
import os
import json

base_path = r'C:\Users\ahmet\source\repos\ScretCustomer\Backend\SecretCustomer.API\wwwroot\ScreenShots\raporlar'

output = []

# List all files
files = os.listdir(base_path)
output.append(f"Klasordeki dosyalar: {files}")

for filename in files:
    filepath = os.path.join(base_path, filename)
    output.append("\n" + "="*80)
    output.append(f"DOSYA: {filename}")
    output.append("="*80)

    try:
        if filename.endswith('.xlsx'):
            wb = openpyxl.load_workbook(filepath)
            for sheet_name in wb.sheetnames:
                output.append(f"\n--- Sheet: {sheet_name} ---")
                sheet = wb[sheet_name]
                row_count = 0
                for row in sheet.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        cells = [str(cell)[:60] if cell else '' for cell in row]
                        output.append(str(cells))
                        row_count += 1
                        if row_count > 100:
                            output.append("... (daha fazla satir var)")
                            break
            wb.close()

        elif filename.endswith('.docx'):
            doc = Document(filepath)
            for para in doc.paragraphs:
                if para.text.strip():
                    output.append(para.text)
            # Tables
            for table in doc.tables:
                output.append("\n--- Tablo ---")
                for row in table.rows:
                    cells = [cell.text.strip()[:40] for cell in row.cells]
                    output.append(str(cells))

        elif filename.endswith('.pptx'):
            prs = Presentation(filepath)
            for i, slide in enumerate(prs.slides):
                output.append(f"\n--- Slide {i+1} ---")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        output.append(shape.text[:200])

    except Exception as e:
        output.append(f"HATA: {str(e)}")

# Write to file
with open(r'C:\Users\ahmet\source\repos\ScretCustomer\reports_output.txt', 'w', encoding='utf-8') as f:
    f.write('\n'.join(output))

print("Rapor okundu ve reports_output.txt dosyasina yazildi.")
