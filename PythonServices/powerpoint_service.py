"""
PowerPoint Report Generation Service for Secret Customer Evaluation System

This service handles:
1. Automated dashboard report generation
2. Branch comparison presentations
3. Trend analysis presentations
4. Custom evaluation reports
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.dml.color import RGBColor
from typing import List, Dict, Any
from datetime import datetime
import json


class PowerPointService:
    """Service for generating PowerPoint presentations"""

    def __init__(self, template_path: str = None):
        """
        Initialize PowerPoint service

        Args:
            template_path: Path to custom PowerPoint template (optional)
        """
        if template_path:
            self.prs = Presentation(template_path)
        else:
            self.prs = Presentation()
            self.prs.slide_width = Inches(10)
            self.prs.slide_height = Inches(7.5)

        # Color scheme
        self.primary_color = RGBColor(68, 114, 196)  # Blue
        self.secondary_color = RGBColor(112, 173, 71)  # Green
        self.accent_color = RGBColor(237, 125, 49)  # Orange
        self.danger_color = RGBColor(192, 0, 0)  # Red
        self.text_color = RGBColor(0, 0, 0)  # Black

    # ==================== SLIDE CREATION ====================

    def create_title_slide(
        self,
        title: str,
        subtitle: str,
        generated_date: str = None
    ):
        """Create title slide"""
        slide_layout = self.prs.slide_layouts[0]  # Title layout
        slide = self.prs.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = self.primary_color

        subtitle_shape = slide.placeholders[1]
        if generated_date:
            subtitle_shape.text = f"{subtitle}\n\nGenerated: {generated_date}"
        else:
            subtitle_shape.text = subtitle

        subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)

    def create_overview_slide(
        self,
        metrics: Dict[str, Any]
    ):
        """Create dashboard overview slide with key metrics"""
        slide_layout = self.prs.slide_layouts[5]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "Dashboard Overview"
        title_frame.paragraphs[0].font.size = Pt(36)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.primary_color

        # Metric boxes (4 columns x 2 rows)
        metrics_list = [
            ("Total Projects", metrics.get('totalProjects', 0), self.primary_color),
            ("Active Projects", metrics.get('activeProjects', 0), self.secondary_color),
            ("Total Evaluations", metrics.get('totalEvaluations', 0), self.accent_color),
            ("Completed", metrics.get('completedEvaluations', 0), self.secondary_color),
            ("Pending", metrics.get('pendingEvaluations', 0), self.danger_color),
            ("Average Score", f"{metrics.get('averageScore', 0):.1f}%", self.primary_color),
            ("Total Branches", metrics.get('totalBranches', 0), self.accent_color),
            ("Total Evaluators", metrics.get('totalEvaluators', 0), self.primary_color)
        ]

        left_positions = [0.5, 2.8, 5.1, 7.4]
        top_positions = [1.5, 4.0]

        for idx, (label, value, color) in enumerate(metrics_list):
            col = idx % 4
            row = idx // 4

            left = Inches(left_positions[col])
            top = Inches(top_positions[row])
            width = Inches(2.0)
            height = Inches(2.0)

            # Background shape
            shape = slide.shapes.add_shape(
                1,  # Rectangle
                left, top, width, height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
            shape.line.color.rgb = color

            # Value text
            text_box = slide.shapes.add_textbox(left, top + Inches(0.3), width, Inches(1.0))
            text_frame = text_box.text_frame
            text_frame.text = str(value)
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            text_frame.paragraphs[0].font.size = Pt(48)
            text_frame.paragraphs[0].font.bold = True
            text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

            # Label text
            label_box = slide.shapes.add_textbox(left, top + Inches(1.3), width, Inches(0.5))
            label_frame = label_box.text_frame
            label_frame.text = label
            label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            label_frame.paragraphs[0].font.size = Pt(14)
            label_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    def create_branch_comparison_slide(
        self,
        branches: List[Dict[str, Any]]
    ):
        """Create branch comparison chart slide"""
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "Branch Performance Comparison"
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.primary_color

        # Create chart
        chart_data = CategoryChartData()
        chart_data.categories = [b.get('branchName', f"Branch {i+1}") for i, b in enumerate(branches)]

        # Add series
        avg_scores = [b.get('averageScore', 0) for b in branches]
        chart_data.add_series('Average Score (%)', avg_scores)

        # Add chart to slide
        x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(5)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        ).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False

        # Chart formatting
        chart.chart_title.text_frame.text = ""
        chart.value_axis.maximum_scale = 100

        # Color the bars
        series = chart.series[0]
        for idx, point in enumerate(series.points):
            fill = point.format.fill
            fill.solid()
            # Green for high scores, orange for medium, red for low
            score = avg_scores[idx]
            if score >= 80:
                fill.fore_color.rgb = self.secondary_color
            elif score >= 60:
                fill.fore_color.rgb = self.accent_color
            else:
                fill.fore_color.rgb = self.danger_color

    def create_trend_analysis_slide(
        self,
        trend_data: List[Dict[str, Any]]
    ):
        """Create trend analysis line chart"""
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "Performance Trend Analysis"
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.primary_color

        # Create chart
        chart_data = CategoryChartData()
        chart_data.categories = [d.get('period', '') for d in trend_data]

        # Add series
        avg_scores = [d.get('averageScore', 0) for d in trend_data]
        completion_rates = [d.get('completionRate', 0) for d in trend_data]

        chart_data.add_series('Average Score (%)', avg_scores)
        chart_data.add_series('Completion Rate (%)', completion_rates)

        # Add chart to slide
        x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(5)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
        ).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.value_axis.maximum_scale = 100

        # Chart formatting
        chart.chart_title.text_frame.text = ""

    def create_evaluator_performance_slide(
        self,
        evaluators: List[Dict[str, Any]]
    ):
        """Create evaluator performance table slide"""
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "Evaluator Performance"
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.primary_color

        # Create table
        rows = len(evaluators) + 1  # +1 for header
        cols = 5
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(4.5)

        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Set column widths
        table.columns[0].width = Inches(2.5)
        table.columns[1].width = Inches(1.5)
        table.columns[2].width = Inches(1.5)
        table.columns[3].width = Inches(1.5)
        table.columns[4].width = Inches(1.0)

        # Header row
        headers = ['Evaluator Name', 'Completed', 'Pending', 'Total', 'Rate']
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.primary_color
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Data rows
        for row_idx, evaluator in enumerate(evaluators, start=1):
            completed = evaluator.get('completedCount', 0)
            pending = evaluator.get('pendingCount', 0)
            total = evaluator.get('totalAssigned', 0)
            rate = (completed / total * 100) if total > 0 else 0

            data = [
                evaluator.get('evaluatorName', ''),
                str(completed),
                str(pending),
                str(total),
                f"{rate:.1f}%"
            ]

            for col_idx, value in enumerate(data):
                cell = table.cell(row_idx, col_idx)
                cell.text = value
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

                # Color code the rate
                if col_idx == 4:
                    if rate >= 80:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(198, 239, 206)
                    elif rate < 60:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(255, 199, 206)

    def create_summary_slide(
        self,
        summary_text: str,
        recommendations: List[str] = None
    ):
        """Create summary and recommendations slide"""
        slide_layout = self.prs.slide_layouts[1]  # Title and content
        slide = self.prs.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = "Summary & Recommendations"

        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()

        # Summary
        p = text_frame.paragraphs[0]
        p.text = summary_text
        p.level = 0
        p.font.size = Pt(18)

        # Recommendations
        if recommendations:
            text_frame.add_paragraph()
            p = text_frame.add_paragraph()
            p.text = "Recommendations:"
            p.font.bold = True
            p.font.size = Pt(20)
            p.space_before = Pt(12)

            for recommendation in recommendations:
                p = text_frame.add_paragraph()
                p.text = recommendation
                p.level = 1
                p.font.size = Pt(16)

    # ==================== REPORT GENERATION ====================

    def generate_dashboard_report(
        self,
        dashboard_data: Dict[str, Any],
        output_path: str,
        project_name: str = "Dashboard Report"
    ) -> str:
        """
        Generate complete dashboard report presentation

        Args:
            dashboard_data: Dictionary containing all dashboard data
            output_path: Path to save the presentation
            project_name: Name of the project/report

        Returns:
            Path to the created presentation
        """
        try:
            # Title slide
            self.create_title_slide(
                title=project_name,
                subtitle="Performance Analysis Report",
                generated_date=datetime.now().strftime('%Y-%m-%d %H:%M')
            )

            # Overview slide
            if 'overview' in dashboard_data:
                self.create_overview_slide(dashboard_data['overview'])

            # Branch comparison
            if 'branchPerformance' in dashboard_data and dashboard_data['branchPerformance']:
                self.create_branch_comparison_slide(dashboard_data['branchPerformance'])

            # Trend analysis
            if 'trendData' in dashboard_data and dashboard_data['trendData']:
                self.create_trend_analysis_slide(dashboard_data['trendData'])

            # Evaluator performance
            if 'evaluatorPerformance' in dashboard_data and dashboard_data['evaluatorPerformance']:
                self.create_evaluator_performance_slide(dashboard_data['evaluatorPerformance'])

            # Summary slide
            summary = self._generate_summary(dashboard_data)
            recommendations = self._generate_recommendations(dashboard_data)
            self.create_summary_slide(summary, recommendations)

            # Save presentation
            self.prs.save(output_path)
            return output_path

        except Exception as e:
            raise Exception(f"Failed to generate dashboard report: {str(e)}")

    def generate_branch_report(
        self,
        branch_data: Dict[str, Any],
        output_path: str
    ) -> str:
        """Generate branch-specific performance report"""
        try:
            branch_name = branch_data.get('branchName', 'Branch')

            # Title slide
            self.create_title_slide(
                title=f"{branch_name} Performance Report",
                subtitle="Detailed Analysis",
                generated_date=datetime.now().strftime('%Y-%m-%d')
            )

            # Branch metrics
            metrics = {
                'totalEvaluations': branch_data.get('totalEvaluations', 0),
                'completedEvaluations': branch_data.get('completed', 0),
                'pendingEvaluations': branch_data.get('pending', 0),
                'averageScore': branch_data.get('averageScore', 0),
                'maxScore': branch_data.get('maxScore', 0),
                'minScore': branch_data.get('minScore', 0)
            }
            self.create_overview_slide(metrics)

            # Trend if available
            if 'trend' in branch_data:
                self.create_trend_analysis_slide(branch_data['trend'])

            # Save
            self.prs.save(output_path)
            return output_path

        except Exception as e:
            raise Exception(f"Failed to generate branch report: {str(e)}")

    # ==================== HELPER METHODS ====================

    def _generate_summary(self, data: Dict[str, Any]) -> str:
        """Generate summary text from dashboard data"""
        overview = data.get('overview', {})
        avg_score = overview.get('averageScore', 0)
        total_evals = overview.get('totalEvaluations', 0)
        completed = overview.get('completedEvaluations', 0)

        completion_rate = (completed / total_evals * 100) if total_evals > 0 else 0

        summary = (
            f"Overall performance shows an average score of {avg_score:.1f}% "
            f"across {total_evals} total evaluations. "
            f"The completion rate stands at {completion_rate:.1f}% with "
            f"{completed} completed evaluations."
        )
        return summary

    def _generate_recommendations(self, data: Dict[str, Any]) -> List[str]:
        """Generate recommendations based on data analysis"""
        recommendations = []

        overview = data.get('overview', {})
        avg_score = overview.get('averageScore', 0)

        if avg_score < 70:
            recommendations.append("Focus on improving overall service quality - current average is below target")

        # Check branch performance
        branches = data.get('branchPerformance', [])
        if branches:
            low_performers = [b for b in branches if b.get('averageScore', 0) < 60]
            if low_performers:
                branch_names = ', '.join([b.get('branchName', '') for b in low_performers[:3]])
                recommendations.append(f"Provide additional training for low-performing branches: {branch_names}")

        # Check completion rate
        completion_rate = overview.get('completedEvaluations', 0) / overview.get('totalEvaluations', 1) * 100
        if completion_rate < 80:
            recommendations.append("Improve evaluation completion rate through better follow-up and reminders")

        if not recommendations:
            recommendations.append("Continue current practices to maintain high performance standards")

        return recommendations


# ==================== STANDALONE USAGE ====================

if __name__ == "__main__":
    import sys

    if len(sys.argv) < 3:
        print("Usage: python powerpoint_service.py <command> <data_json> <output_path>")
        print("\nCommands:")
        print("  dashboard <data_json> <output_path>")
        print("  branch <data_json> <output_path>")
        sys.exit(1)

    command = sys.argv[1]
    data_path = sys.argv[2]
    output_path = sys.argv[3] if len(sys.argv) > 3 else "report.pptx"

    try:
        with open(data_path, 'r', encoding='utf-8') as f:
            data = json.load(f)

        service = PowerPointService()

        if command == "dashboard":
            result = service.generate_dashboard_report(data, output_path)
            print(f"Dashboard report created: {result}")

        elif command == "branch":
            result = service.generate_branch_report(data, output_path)
            print(f"Branch report created: {result}")

        else:
            print(f"Unknown command: {command}")
            sys.exit(1)

    except Exception as e:
        print(f"Error: {str(e)}", file=sys.stderr)
        sys.exit(1)
